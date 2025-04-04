
import streamlit as st
import requests
from fpdf import FPDF
from docx import Document
import openai
import PyPDF2
from docx.shared import Inches
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import zipfile
import os
from pathlib import Path
import csv

openai.api_key = st.secrets["api"]["OPENAI_API_KEY"]


def generate_detailed_ppt_content(domain, topic):
    """Generate detailed content for a presentation using GPT based on the selected domain and topic."""
    prompt = (
        f"You are an expert in the {domain} domain. Generate a professional, formal PowerPoint presentation on the topic: '{topic}'.\n\n"
        f"Instructions:\n"
        f"1. All content must be specific to the {domain} domain and based on the topic '{topic}'.\n"
        f"2. Each slide must have at least 4 well-written bullet points, not paragraphs.\n"
        f"3. Ensure formal tone, clarity, and relevance to the chosen domain.\n"
        f"4. Structure should include:\n"
        f"   - Title Slide (Domain, Topic, Author Name placeholder)\n"
        f"   - Introduction Slide (Definition and importance of the topic)\n"
        f"   - 4–6 Key Point Slides (with elaborated points)\n"
        f"   - Case Studies/Examples Slide (with real-world relevance to the domain)\n"
        f"   - Conclusion Slide (with summary/future direction)\n\n"
        f"Output only detailed slide-wise content."
    )
    try:
        response = openai.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": f"You are a domain expert in {domain}."},
                {"role": "user", "content": prompt},
            ],
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"

def create_professional_ppt(content, topic, file_name="presentation.pptx"):
    """Create a well-formatted professional PowerPoint presentation."""
    ppt = Presentation()

    # Set consistent font styles
    def set_textbox_style(text_frame):
        """Style the textbox content."""
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = "Calibri (Body)"
            paragraph.font.size = Pt(20)
            paragraph.alignment = PP_ALIGN.LEFT

    # Title Slide
    title_slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = "A Comprehensive Overview in {domain} Domain"

    # Process GPT content into slides
    sections = content.split("\n\n")
    for section in sections:
        if ":" in section:  # Detect title: content structure
            slide_title, slide_content = section.split(":", 1)

            # Add a new slide for each section
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_title.strip()

            # Add content to the slide
            content_box = slide.placeholders[1]
            content_box.text = slide_content.strip()
            set_textbox_style(content_box.text_frame)

    # Save the presentation
    ppt.save(file_name)
    return file_name



# Set up the page configuration (must be the first command)
st.set_page_config(page_title="AI-Powered Content Generation", layout="wide", page_icon="📚")

# Title Section with enhanced visuals
st.markdown(
    """
    <h1 style="text-align: center; font-size: 2.5rem; color: #4A90E2;">📚 AI-Powered PPT Content Generation</h1>
    <p style="text-align: center; font-size: 1.1rem; color: #555;">Streamline your content creation process with AI technology. Designed for the <strong>pharmaceutical</strong> and <strong>medical</strong> domains.</p>
    """,
    unsafe_allow_html=True,
)
# Horizontal line
st.markdown("---")

# Content Generation Instructions
with st.expander("1️⃣ **PPT Generation Instructions**", expanded=True):
    st.markdown("""
        - Create professional PowerPoint presentations based on medical and pharmaceutical topics.
        - **Steps**:
          1. Enter a valid topic in the input field.
          2. Click the **Generate PPT** button to create a detailed presentation.
          3. Download the PPT file using the download button provided.
        """)

# Horizontal line
st.markdown("---")

st.header("📊 PPT Content Generation")

# Step 1: Get the domain
domain = st.text_input("Enter the domain for your presentation:", placeholder="e.g., Medical, Finance, Education")

# Step 2: Get the topic
topic = ""
if domain:
    topic = st.text_input(f"Enter the topic related to the {domain} domain:", placeholder="e.g., Drug Discovery, Stock Market Trends, Online Learning Platforms")

# Step 3: Generate and download PPT
if st.button("Generate PPT"):
    if domain and topic:
        st.info("Generating detailed content for your presentation. Please wait...")
        detailed_content = generate_detailed_ppt_content(domain, topic)
        if "Error" not in detailed_content:
            ppt_file_name = create_professional_ppt(detailed_content, f"{domain} - {topic}")
            st.success("Your PowerPoint presentation has been successfully generated!")
            with open(ppt_file_name, "rb") as file:
                st.download_button(
                    "Download Your PPT",
                    file,
                    file_name=ppt_file_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
        else:
            st.error(detailed_content)
    else:
        st.warning("Please enter both domain and topic before generating the presentation.")

# Horizontal line
st.markdown("---")

# Footer
st.caption("Developed by **Corbin Technology Solutions**")
